import sys
from PyQt5.QtWidgets import  QApplication, QWidget, QVBoxLayout, QPushButton, QProgressBar, QLabel, QFileDialog, QSpinBox, QTabWidget, QHBoxLayout
from PyQt5.QtCore import QThread, pyqtSignal, Qt
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

class ConvertThread(QThread):
    error = pyqtSignal(str)
    progress = pyqtSignal(int)
    finished = pyqtSignal()

    def __init__(self, pdf_path, pptx_path, dpi):
        super().__init__()
        self.pdf_path = pdf_path
        self.pptx_path = pptx_path
        self.dpi = dpi

    def run(self):
        try:
            doc = fitz.open(self.pdf_path)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(10 * 9 / 16)
            total_pages = len(doc)
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=self.dpi)
                img_data = pix.tobytes("PNG")
                img_stream = BytesIO(img_data)
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=prs.slide_width)
                img_stream.close()
                self.progress.emit(int((i + 1) / total_pages * 100))
            prs.save(self.pptx_path)
        except Exception as e:
            self.error.emit(str(e))
        finally:
            self.finished.emit()

class PDFtoPPTXConverter(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()
        # 禁用最大化按钮
        self.setWindowFlags(self.windowFlags() & ~Qt.WindowMaximizeButtonHint)
        # 设置固定窗口大小
        self.setFixedSize(440, 250)
        self.pdf_path = None  # 初始化pdf_path属性
        self.pptx_path = None  # 初始化pptx_path属性


    def initUI(self):
        self.setWindowTitle("PDF to PPTX Converter")
        self.setGeometry(100, 100, 400, 300)  # 调整窗口大小以适应选项卡

        # 创建主布局
        mainLayout = QVBoxLayout()
        self.setLayout(mainLayout)

        # 创建选项卡控件
        tabWidget = QTabWidget()
        mainLayout.addWidget(tabWidget)

        # 创建主界面选项卡
        mainTab = QWidget()
        tabWidget.addTab(mainTab, "Main")
        mainTabLayout = QVBoxLayout()
        mainTab.setLayout(mainTabLayout)

        # 创建设置选项卡
        settingsTab = QWidget()
        tabWidget.addTab(settingsTab, "Settings")
        settingsLayout = QVBoxLayout()
        settingsTab.setLayout(settingsLayout)

        # 在设置选项卡中添加DPI设置及其标签
        dpiLabel = QLabel("DPI Setting (Range: 72 to 600):")
        settingsLayout.addWidget(dpiLabel)

        self.dpiSpinBox = QSpinBox()
        self.dpiSpinBox.setRange(72, 600)  # 设置DPI的范围为72到600
        self.dpiSpinBox.setValue(300)  # 默认值设置为300
        settingsLayout.addWidget(self.dpiSpinBox)

        # 使用固定大小的容器来固定布局
        settingsContainer = QWidget()
        settingsContainer.setLayout(settingsLayout)
        settingsContainer.setFixedSize(380, 100)  # 设置固定大小

        # 将容器添加到选项卡中而不是直接添加布局
        settingsTabLayout = QVBoxLayout()
        settingsTabLayout.addWidget(settingsContainer)
        settingsTab.setLayout(settingsTabLayout)

        # 对主界面选项卡进行相同的处理
        mainTabContainer = QWidget()
        mainTabContainer.setLayout(mainTabLayout)
        mainTabContainer.setFixedSize(380, 180)  # 根据需要调整大小

        mainTabLayoutContainer = QVBoxLayout()
        mainTabLayoutContainer.addWidget(mainTabContainer)
        mainTab.setLayout(mainTabLayoutContainer)

        # 在主界面选项卡中添加控件
        self.btnSelectPDF = QPushButton('Select PDF', self)
        self.btnSelectPDF.clicked.connect(self.selectPDF)
        mainTabLayout.addWidget(self.btnSelectPDF)

        self.btnConvert = QPushButton('Convert', self)
        self.btnConvert.clicked.connect(self.startConversion)
        mainTabLayout.addWidget(self.btnConvert)

        self.progressBar = QProgressBar(self)
        self.progressBar.setAlignment(Qt.AlignCenter)
        mainTabLayout.addWidget(self.progressBar)

        self.infoLabel = QLabel('', self)  # 通用信息显示标签
        self.infoLabel.setStyleSheet("color: black;")  # 默认信息显示为黑色
        mainTabLayout.addWidget(self.infoLabel)

    def selectPDF(self):
        self.pdf_path, _ = QFileDialog.getOpenFileName(self, "Select PDF file", "", "PDF files (*.pdf)")
        if not self.pdf_path:
            self.infoLabel.setStyleSheet("color: red;")
            self.infoLabel.setText("Warning: No PDF file selected.")
            return
        self.pptx_path = self.pdf_path.replace('.pdf', '.pptx')

    def startConversion(self):
        if not self.pdf_path or not self.pptx_path:  # 检查是否已选择PDF文件
            self.infoLabel.setStyleSheet("color: red;")
            self.infoLabel.setText("Warning: Please select a PDF file first.")
            return        
        dpi_value = self.dpiSpinBox.value()  # 获取用户设置的DPI值
        if dpi_value < 72 or dpi_value > 600:
            self.infoLabel.setStyleSheet("color: red;")
            self.infoLabel.setText("Warning: DPI value must be between 72 and 600.")
            return
        self.thread = ConvertThread(self.pdf_path, self.pptx_path, dpi_value)
        self.thread.error.connect(self.showError)
        self.thread.progress.connect(self.updateProgress)
        self.thread.finished.connect(self.conversionFinished)
        self.thread.start()

    def showError(self, message):
        self.infoLabel.setStyleSheet("color: red;")  # 错误信息显示为红色
        self.infoLabel.setText(message)  # 显示错误信息   

    def updateProgress(self, value):
        self.progressBar.setValue(value)
        self.infoLabel.setText(f'{value}%')
        if value < 100:
            self.infoLabel.setStyleSheet("color: black;")  # 进度信息显示为黑色
            self.infoLabel.setText("Converting...")
        else:
            self.conversionFinished()

    def conversionFinished(self):
        self.infoLabel.setStyleSheet("color: green;")  # 完成信息显示为绿色
        self.infoLabel.setText("Conversion completed successfully.")  # 显示转换完成信息

if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = PDFtoPPTXConverter()
    ex.show()
    sys.exit(app.exec_())